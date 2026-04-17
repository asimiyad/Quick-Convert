import os
from engine.core.base_engine import BaseConverter, validate_file
from engine.config import SUPPORTED_CONVERSIONS

try:
    import docx # python-docx
except ImportError:
    docx = None

class WordConverter(BaseConverter):
    """
    Handles converting Word files (.docx, .doc) into other formats (TXT, PDF, PPTX).
    Optimized for headless cloud environments (Render) - No LibreOffice required.
    """

    def validate_requirements(self) -> bool:
        """
        Validates input payload and configuration mappings safely.
        """
        # 1. Base validation (Existence & Signature check)
        if not validate_file(self.input_file):
            return False

        # 2. Config mapped support check
        _, input_ext = os.path.splitext(self.input_file)
        _, output_ext = os.path.splitext(self.output_file)
        input_ext, output_ext = input_ext.lower(), output_ext.lower()
        supported_outputs = SUPPORTED_CONVERSIONS.get(input_ext, [])
        
        if output_ext not in supported_outputs:
            self.logger.error(f"Conversion from {input_ext} to {output_ext} is disabled/unsupported by config.")
            return False

        # 3. Dynamic Dependency Checks
        if output_ext in [".txt", ".pptx"] and docx is None:
            self.logger.error("Missing Python module 'python-docx'. Aborting.")
            return False

        self.logger.info("Word format validation successful.")
        return True

    def convert(self) -> bool:
        """
        Root executor for Word logic.
        """
        _, output_ext = os.path.splitext(self.output_file)
        output_ext = output_ext.lower()
        
        try:
            if output_ext == ".txt":
                return self._convert_to_txt()
            elif output_ext == ".pdf":
                return self._convert_to_pdf()
            elif output_ext == ".pptx":
                return self._convert_to_pptx()
            else:
                self.logger.error(f"Worker logic for {output_ext} is not yet implemented in WordConverter.")
                return False
        except Exception as e:
            self.logger.exception(f"Critical error during physical Word conversion: {e}")
            return False

    def _convert_to_txt(self) -> bool:
        """
        Extracts structural text elements natively using python-docx.
        """
        self.logger.debug(f"Beginning text extraction sequence for: {self.input_file}")
        try:
            document = docx.Document(self.input_file)
            text_content = [para.text for para in document.paragraphs if para.text.strip()]
            
            with open(self.output_file, 'w', encoding='utf-8') as text_file:
                text_file.write("\n\n".join(text_content))
                
            self.logger.info(f"Successfully generated clean TXT: {self.output_file}")
            return True
        except Exception as e:
            self.logger.error(f"Physical text Extraction halted: {e}")
            return False

    def _convert_to_pdf(self) -> bool:
        """
        Converts DOCX to PDF prioritizing docx2pdf (Windows Native), fallback to LibreOffice (Linux),
        and finally pure-Python Aspose.Words.
        """
        import platform
        import subprocess
        import shutil

        self.logger.info("Engaging multi-fallback DOCX -> PDF mapping.")

        # 1. Try docx2pdf natively (Perfect for Windows w/ MS Word)
        if platform.system() == "Windows":
            try:
                from docx2pdf import convert as docx_convert
                self.logger.info("Windows detected: Attempting docx2pdf payload routing.")
                docx_convert(self.input_file, self.output_file)
                self.logger.info(f"Successfully generated clean PDF via docx2pdf: {self.output_file}")
                return True
            except Exception as e:
                self.logger.warning(f"docx2pdf routing failed or MS Word not natively installed: {e}")

        # 2. Try LibreOffice Headless (Perfect for Linux/Render)
        self.logger.info("Attempting fallback sequence to native LibreOffice headless engine...")
        outdir = os.path.dirname(self.output_file)
        command = [
            "soffice", "--headless", "--nologo", "--nofirststartwizard",
            "--convert-to", "pdf", self.input_file, "--outdir", outdir
        ]
        try:
            result = subprocess.run(command, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
            if result.returncode == 0:
                base_filename = os.path.basename(self.input_file)
                name_no_ext, _ = os.path.splitext(base_filename)
                lo_output_file = os.path.join(outdir, f"{name_no_ext}.pdf")
                if lo_output_file != self.output_file and os.path.exists(lo_output_file):
                    shutil.move(lo_output_file, self.output_file)
                self.logger.info(f"Successfully generated native PDF via LibreOffice: {self.output_file}")
                return True
            else:
                 self.logger.warning(f"LibreOffice backend failed natively: {result.stderr}")
        except Exception as e:
            self.logger.warning(f"LibreOffice command totally unresolvable natively: {e}")

        # 3. Final Fallback: Aspose.Words (Will apply an evaluation watermark)
        self.logger.info("Executing final structural fallback sequence via pure-Python Aspose.Words engine.")
        try:
            import aspose.words as aw
            doc = aw.Document(self.input_file)
            doc.save(self.output_file)
            self.logger.warning(f"Generated PDF via Aspose engine (Evaluation Watermark inevitably applied): {self.output_file}")
            return True
        except Exception as e:
            self.logger.error(f"Aspose rendering severely physically halted: {e}")
            return False

    def _convert_to_pptx(self) -> bool:
        """
        Constructs PPTX natively from DOCX formatting.
        """
        self.logger.info("Automatically architecting PPTX slide deck natively from DOCX formatting...")
        try:
            import pptx
            from pptx.util import Pt
            import docx as python_docx
            
            doc = python_docx.Document(self.input_file)
            prs = pptx.Presentation()
            
            for i in range(0, len(doc.paragraphs), 6):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = slide.shapes.title
                title.text = f"Content Section {i//6 + 1}"
                
                tf = slide.placeholders[1].text_frame
                chunk = doc.paragraphs[i:i+6]
                for para in chunk:
                    if para.text.strip():
                        p = tf.add_paragraph()
                        p.text = para.text.strip()
                        p.font.size = Pt(16)
                        
            prs.save(self.output_file)
            self.logger.info(f"Presentation securely compiled successfully: {self.output_file}")
            return True
        except Exception as e:
            self.logger.error(f"Docx bridging halted explicitly physically: {e}")
            return False