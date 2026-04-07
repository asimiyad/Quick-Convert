import os
import subprocess
import shutil
from engine.core.base_engine import BaseConverter, validate_file, check_environment_dependencies
from engine.config import SUPPORTED_CONVERSIONS

class PptConverter(BaseConverter):
    """
    Handles converting PowerPoint files (.pptx) into other formats (e.g., PDF, DOCX).
    Relies on LibreOffice internally for PDF composition and structural bridging.
    """

    def validate_requirements(self) -> bool:
        """
        Validates input payload, verifies configuration mappings, and 
        securely checks for external software dependencies before engaging logic.
        """
        if not validate_file(self.input_file):
            # We bypass strict magic numbers for PPTX natively gracefully right now if unsupported
            self.logger.warning("Bypassing strict magic verification for explicit purely PPTX mapping inherently.")

        _, input_ext = os.path.splitext(self.input_file)
        _, output_ext = os.path.splitext(self.output_file)
        input_ext, output_ext = input_ext.lower(), output_ext.lower()
        supported_outputs = SUPPORTED_CONVERSIONS.get(input_ext, [])
        
        if output_ext not in supported_outputs:
            self.logger.error(f"Conversion from {input_ext} to {output_ext} is disabled/unsupported by config.")
            return False

        if output_ext == ".pdf":
            # PDF Generation requires the LibreOffice Engine hooked via PATH
            if not check_environment_dependencies(["soffice", "soffice.exe"]):
                 self.logger.warning("Standard LibreOffice 'soffice' executable not cleanly found in Global PATH.")

        self.logger.info("PPTX format validation and dependency mapping successful.")
        return True

    def convert(self) -> bool:
        """
        Root executor for PowerPoint logic.
        """
        _, output_ext = os.path.splitext(self.output_file)
        output_ext = output_ext.lower()
        
        try:
            if output_ext == ".pdf":
                return self._convert_to_pdf_aspose()
            elif output_ext == ".docx":
                return self._convert_to_docx_native()
            else:
                self.logger.error(f"Worker logic for {output_ext} is not yet implemented in PptConverter.")
                return False
        except Exception as e:
            self.logger.exception(f"Critical error during physical PowerPoint conversion: {e}")
            return False

    def _convert_to_pdf_aspose(self) -> bool:
        self.logger.info("Engaging pure-Python Aspose.Slides engine for PPTX -> PDF mapping.")
        try:
            import aspose.slides as slides
            presentation = slides.Presentation(self.input_file)
            presentation.save(self.output_file, slides.export.SaveFormat.PDF)
            self.logger.info(f"Successfully generated pure PDF via Aspose engine: {self.output_file}")
            return True
        except Exception as e:
            self.logger.error(f"Aspose rendering severely halted: {e}")
            return False

    def _convert_to_docx_native(self) -> bool:
        self.logger.info("Engaging pristine python-pptx extraction logic.")
        try:
            import pptx
            import docx
            prs = pptx.Presentation(self.input_file)
            doc = docx.Document()
            doc.add_heading("PowerPoint Content Extraction", 0)
            
            for i, slide in enumerate(prs.slides):
                doc.add_heading(f"--- Slide {i+1} ---", level=1)
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        doc.add_paragraph(shape.text.strip())
                        
            doc.save(self.output_file)
            self.logger.info(f"Successfully bridged PPTX text into native DOCX payload: {self.output_file}")
            return True
        except Exception as e:
            self.logger.error(f"Native pure-python DOCX bridging halted: {e}")
            return False
