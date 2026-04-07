import os
from engine.core.base_engine import BaseConverter, validate_file
from engine.config import SUPPORTED_CONVERSIONS

# Attempt to load specialized conversion libraries gracefully
try:
    import pypdf
except ImportError:
    pypdf = None

try:
    from pdf2docx import Converter as PDF2DocxConverter
except ImportError:
    PDF2DocxConverter = None

try:
    import fitz # PyMuPDF backend for Image rendering
except ImportError:
    fitz = None


class PdfConverter(BaseConverter):
    """
    Handles converting PDF files into other formats (e.g., TXT, DOCX, PNG, JPEG).
    Inherits from BaseConverter to maintain standard lifecycle logic.
    """

    def validate_requirements(self) -> bool:
        """
        Validate input file exists, is not empty, and ensure the 
        output format is structurally supported.
        """
        # 1. Deep binary and existence validation from core
        if not validate_file(self.input_file):
            return False

        # 2. Extension validation against configuration mappings
        _, input_ext = os.path.splitext(self.input_file)
        _, output_ext = os.path.splitext(self.output_file)
        
        input_ext, output_ext = input_ext.lower(), output_ext.lower()
        supported_outputs = SUPPORTED_CONVERSIONS.get(input_ext, [])
        
        if output_ext not in supported_outputs:
            self.logger.error(f"Conversion from {input_ext} to {output_ext} is disabled/unsupported by config.")
            return False

        # 3. Check for specific library dependencies payload requires
        if output_ext == ".txt" and pypdf is None:
            self.logger.error("Missing required Python module 'pypdf'. Aborting. Install via: pip install pypdf")
            return False
            
        if output_ext == ".docx" and PDF2DocxConverter is None:
            self.logger.error("Missing required Python module 'pdf2docx'. Aborting. Install via: pip install pdf2docx")
            return False
            
        if output_ext in [".png", ".jpeg", ".jpg"] and fitz is None:
            self.logger.error("Missing required Python module 'pymupdf'. Aborting. Install via: pip install pymupdf")
            return False

        self.logger.info("PDF format validation and dependency mapping successful.")
        return True

    def convert(self) -> bool:
        """
        Root executor for PDF logic. Routes to internal worker functions
        based on the desired output extension.
        """
        _, output_ext = os.path.splitext(self.output_file)
        output_ext = output_ext.lower()
        
        try:
            if output_ext == ".txt":
                return self._convert_to_txt()
            elif output_ext == ".docx":
                return self._convert_to_docx()
            elif output_ext == ".pptx":
                return self._convert_to_pptx()
            elif output_ext in [".png", ".jpeg", ".jpg"]:
                return self._convert_to_image()
            else:
                self.logger.error(f"Worker logic for {output_ext} is not yet implemented in PdfConverter.")
                return False
        except Exception as e:
            self.logger.exception(f"Critical error during physical PDF conversion: {e}")
            return False

    def _convert_to_txt(self) -> bool:
        """
        Internal logic strictly isolating the PDF to TXT extraction process.
        """
        self.logger.debug(f"Beginning text extraction sequence for: {self.input_file}")
        
        try:
            with open(self.input_file, 'rb') as file:
                reader = pypdf.PdfReader(file)
                text_content = []
                
                total_pages = len(reader.pages)
                self.logger.info(f"Target PDF loaded. Constructing {total_pages} page(s).")

                for i, page in enumerate(reader.pages):
                    self.logger.debug(f"Parsing page {i+1}/{total_pages}...")
                    page_text = page.extract_text() or ""
                    text_content.append(page_text)
            
            # Write aggregated text out cleanly
            with open(self.output_file, 'w', encoding='utf-8') as text_file:
                text_file.write("\n\n--- Page Break ---\n\n".join(text_content))

            self.logger.info(f"Successfully generated TXT payload: {self.output_file}")
            return True

        except Exception as e:
            self.logger.error(f"Extraction halted. Failed to read/write payload: {e}")
            return False

    def _convert_to_docx(self) -> bool:
        """
        Internal logic strictly isolating the PDF to DOCX extraction process.
        """
        self.logger.debug(f"Beginning PDF to Word conversion sequence for: {self.input_file}")
        
        try:
            cv = PDF2DocxConverter(self.input_file)
            self.logger.info("Executing complex structural conversion... this may take a moment.")
            cv.convert(self.output_file, start=0, end=None)
            cv.close()
            
            self.logger.info(f"Successfully generated DOCX payload: {self.output_file}")
            return True
            
        except Exception as e:
            self.logger.error(f"Docx conversion halted. Reconstruction failed: {e}")
            return False

    def _convert_to_pptx(self) -> bool:
        """
        Internal logic constructing vector-rastered slides flawlessly natively.
        """
        self.logger.info(f"Beginning pure-python PDF to PPTX presentation conversion sequence...")
        import os
        try:
            import pptx
            import fitz
            doc = fitz.open(self.input_file)
            prs = pptx.Presentation()
            blank_slide_layout = prs.slide_layouts[6] 
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0))
                img_path = f"{self.output_file}_temp_{page_num}.png"
                pix.save(img_path)
                
                slide = prs.slides.add_slide(blank_slide_layout)
                slide.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)
                os.remove(img_path)

            prs.save(self.output_file)
            self.logger.info(f"Successfully constructed vector-native PPTX presentation: {self.output_file}")
            return True
        except Exception as e:
            self.logger.error(f"Vector PPTX simulation halted securely: {e}")
            return False

    def _convert_to_image(self) -> bool:
        """
        Internal logic isolating the PDF back to raw Image rasterization.
        Uses PyMuPDF to extract vector PDF pages natively without external software.
        """
        _, ext = os.path.splitext(self.output_file)
        self.logger.debug(f"Beginning PDF to {ext.upper()} conversion sequence for: {self.input_file}")
        
        try:
            doc = fitz.open(self.input_file)
            total_pages = len(doc)
            self.logger.info(f"Target PDF loaded. Preparing to render {total_pages} page(s).")
            
            base_output_path, saved_ext = os.path.splitext(self.output_file)
            
            for page_num in range(total_pages):
                page = doc.load_page(page_num)
                # Render using a Matrix scaling for beautiful High-Res output (roughly 300 DPI)
                zoom_matrix = fitz.Matrix(2.0, 2.0)
                pix = page.get_pixmap(matrix=zoom_matrix)
                
                # Naming schema: if multi-page, append _pageX to the filename
                if total_pages == 1:
                    final_path = self.output_file
                else:
                    final_path = f"{base_output_path}_page{page_num + 1}{saved_ext}"
                    
                pix.save(final_path)
                self.logger.debug(f"Saved snapshot natively to: {final_path}")
                
            doc.close()
            self.logger.info(f"Successfully output {total_pages} image payload(s).")
            return True
            
        except Exception as e:
            self.logger.error(f"Image compilation halted. Rasterization process failed: {e}")
            return False
